import sys
from pptx import Presentation
from openpyxl import load_workbook
from datetime import datetime

EXCEL_PATH = "data/data.xlsx"
PPT_TEMPLATE_PATH = "template/template.pptx"


# ---------------------------------------------------------
# Format Excel cell values (dates, None, numbers, strings)
# ---------------------------------------------------------
def format_value(v):
    if isinstance(v, datetime):
        return v.strftime("%Y-%m-%d")  # clean date only
    if v is None:
        return ""
    return str(v)


# ---------------------------------------------------------
# Read one row from Excel and return {header: value} dict
# ---------------------------------------------------------
def read_row_from_excel(path, row_number):
    wb = load_workbook(path, data_only=True)
    ws = wb.active

    headers = [str(h).strip() for h in next(ws.iter_rows(values_only=True))]
    row = list(ws.iter_rows(min_row=row_number, max_row=row_number, values_only=True))[0]

    entry = {h: format_value(v) for h, v in zip(headers, row)}
    return entry


# ---------------------------------------------------------
# Replace all placeholders {{Field}} inside a slide
# Handles broken runs + empty textboxes safely
# ---------------------------------------------------------
def replace_in_runs(slide, mapping):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame

        # Step 1 — combine all text from all runs
        full_text = ""
        for para in tf.paragraphs:
            for run in para.runs:
                full_text += run.text

        # Step 2 — replace placeholders
        for key, value in mapping.items():
            token = f"{{{{{key}}}}}"
            full_text = full_text.replace(token, value)

        # Step 3 — clear all runs
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = ""

        # Step 4 — ensure at least one paragraph + one run exists
        if len(tf.paragraphs) == 0:
            p = tf.add_paragraph()
            run = p.add_run()
        else:
            p = tf.paragraphs[0]
            if len(p.runs) == 0:
                run = p.add_run()
            else:
                run = p.runs[0]

        # Step 5 — write final text
        run.text = full_text


# ---------------------------------------------------------
# Main program
# ---------------------------------------------------------
def main(row_number=2, output_filename=None):
    entry = read_row_from_excel(EXCEL_PATH, row_number)
    prs = Presentation(PPT_TEMPLATE_PATH)

    # Replace placeholders everywhere
    for slide in prs.slides:
        replace_in_runs(slide, entry)

    # Default output name based on first column
    if output_filename is None:
        first_col_value = list(entry.values())[0]
        safe_name = str(first_col_value).strip().replace(" ", "_")
        output_filename = f"output/{safe_name}.pptx"

    prs.save(output_filename)
    print(f"Presentation generated from Excel row {row_number}: {output_filename}")


# ---------------------------------------------------------
# CLI entry
# ---------------------------------------------------------
if __name__ == "__main__":
    row_number = 2
    output_filename = None

    if len(sys.argv) > 1:
        try:
            row_number = int(sys.argv[1])
        except ValueError:
            print("Please provide a valid row number (integer).")
            sys.exit(1)

    if len(sys.argv) > 2:
        output_filename = sys.argv[2]

    main(row_number, output_filename)
